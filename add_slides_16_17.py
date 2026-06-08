
"""
슬라이드 16/17 가격박스 추가
- 슬라이드 16: 21개 단지
- 슬라이드 17: 11개 단지
- 폰트 10pt / 입출력: CLAUDE 폴더 PPT
"""
import re, json, shutil, openpyxl
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

A    = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
SLIDE_W = 9144000
SLIDE_H = 6858000

RAINBOW = {
    35:('000000','FFFFFF'), 34:('1A1A1A','FFFFFF'), 33:('2B2B2B','FFFFFF'),
    32:('424242','FFFFFF'), 31:('616161','FFFFFF'), 30:('757575','FFFFFF'),
    29:('4A0000','FFFFFF'), 28:('6A0F0F','FFFFFF'), 27:('8D1515','FFFFFF'),
    26:('C62828','FFFFFF'), 25:('E53935','FFFFFF'), 24:('F4511E','FFFFFF'),
    23:('FB8C00','FFFFFF'), 22:('FFB300','FFFFFF'), 21:('F9E400','000000'),
    20:('9CCC65','000000'), 19:('43A047','FFFFFF'), 18:('00897B','FFFFFF'),
    17:('00ACC1','000000'), 16:('29B6F6','000000'), 15:('1E88E5','FFFFFF'),
    14:('1565C0','FFFFFF'), 13:('1A237E','FFFFFF'), 12:('7B2FBE','FFFFFF'),
    11:('B07FE0','FFFFFF'), 10:('D7B8F5','000000'),
}

def fix(v):
    if v is None or isinstance(v, str): return None
    return v // 10 if v > 1000000 else v

def num(m):
    if m is None: return '-'
    v = m / 10000
    s = f'{v:.1f}'
    return s[:-2] if s.endswith('.0') else s

def get_color(mae):
    if mae is None: return ('595959','FFFFFF')
    eok = int(mae // 10000)
    if eok >= 35: return RAINBOW[35]
    return RAINBOW.get(eok, ('595959','FFFFFF'))

def yr2(y):
    if not y: return '??'
    yr = re.sub(r'[^\d].*', '', str(y))
    return yr[-2:] if len(yr) >= 2 else yr

def sd_str(s):
    if s is None: return '?'
    return str(int(s)) if isinstance(s, float) else str(s)

def text_width_pt(text, font_pt=10):
    w = 0
    for c in text:
        if '가' <= c <= '힣':    w += font_pt * 1.0
        elif '一' <= c <= '鿿':  w += font_pt * 1.0
        elif c == ' ':           w += font_pt * 0.3
        elif c in '()':          w += font_pt * 0.4
        elif c in '/.,':         w += font_pt * 0.35
        else:                    w += font_pt * 0.6
    return w

def calc_box_w(all_lines, font_pt=10, min_pt=70, padding_pt=14):
    max_w = max(text_width_pt(l, font_pt) for l in all_lines)
    return max(int(Pt(max_w + padding_pt)), int(Pt(min_pt)))

def set_highlight(run, bg_hex, fg_hex):
    rPr = run._r.find(f'{{{A}}}rPr')
    if rPr is None: rPr = etree.SubElement(run._r, f'{{{A}}}rPr')
    for tag in [f'{{{A}}}highlight', f'{{{A}}}solidFill']:
        for el in rPr.findall(tag): rPr.remove(el)
    sf = etree.SubElement(rPr, f'{{{A}}}solidFill')
    etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', fg_hex)
    hl = etree.SubElement(rPr, f'{{{A}}}highlight')
    etree.SubElement(hl, f'{{{A}}}srgbClr').set('val', bg_hex)

def set_no_fill(shape):
    sp = shape._element
    spPr = sp.find(f'{{{P_NS}}}spPr')
    if spPr is None: spPr = sp.find(f'{{{A}}}spPr')
    if spPr is None: return
    for tag in [f'{{{A}}}solidFill', f'{{{A}}}noFill', f'{{{A}}}gradFill']:
        for el in spPr.findall(tag): spPr.remove(el)
    etree.SubElement(spPr, f'{{{A}}}noFill')

def set_border(shape, hex_color='333333', width_pt=1.0):
    sp = shape._element
    spPr = sp.find(f'{{{P_NS}}}spPr')
    if spPr is None: spPr = sp.find(f'{{{A}}}spPr')
    if spPr is None: return
    for old_ln in spPr.findall(f'{{{A}}}ln'): spPr.remove(old_ln)
    ln = etree.SubElement(spPr, f'{{{A}}}ln')
    ln.set('w', str(int(Pt(width_pt))))
    sf = etree.SubElement(ln, f'{{{A}}}solidFill')
    etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', hex_color)

def set_bPr(tf, l=36000, r=36000, t=18000, b=18000, anchor='t'):
    bPr = tf._txBody.find(f'{{{A}}}bodyPr')
    if bPr is not None:
        bPr.set('lIns', str(l)); bPr.set('rIns', str(r))
        bPr.set('tIns', str(t)); bPr.set('bIns', str(b))
        bPr.set('anchor', anchor)

# ── hanpan5.html 추가 정보 ──────────────────────────────────────
def load_extra_info(html_path):
    with open(html_path, 'rb') as f:
        content = f.read().decode('utf-8')
    m = re.search(r'apartments=(\[.*?\]);', content, re.DOTALL)
    if not m: return {}
    arr = json.loads(m.group(1))
    price_pat = re.compile(r'^\d+p\(')
    result = {}
    for apt in arr:
        name_raw = apt.get('name', '').split('\n')[0].strip()
        text = apt.get('text', '')
        lines = text.split('\n')
        extra = [l.strip() for l in lines[1:] if l.strip() and not price_pat.match(l.strip())]
        # 연식/세대수 패턴 필터링
        filtered = [l for l in extra if not re.match(r"^\d{2}'\s*[\d,]+\^", l)]
        result[name_raw] = filtered  # 빈 리스트도 저장 (매칭용)
    return result

extra_db = load_extra_info('C:/Temp/hanpan5.html')

def get_extra_lines(apt_name):
    """단지명 앞부분으로 hanpan5 매칭 (길이 순 내림차순 시도)"""
    best_match = None
    best_len = 0
    for html_name, lines in extra_db.items():
        # apt_name의 앞 N글자가 html_name에 포함되는지 확인
        for n in range(min(10, len(apt_name)), 3, -1):
            prefix = apt_name[:n]
            if prefix in html_name:
                if n > best_len:
                    best_len = n
                    best_match = lines
                break
    return best_match if best_match else []

# ── Excel 로드 ─────────────────────────────────────────────────
S16_CODES = [629, 631, 626, 627, 628, 147185, 114947, 118141, 115869,
             632, 633, 617, 619, 618, 383, 3431, 103139, 615, 22809, 612]
S17_CODES = [598, 593, 596, 597, 600, 3033, 609, 608, 107544, 166537, 365]
ALL_CODES = list(set(S16_CODES + S17_CODES))

wb = openpyxl.load_workbook('C:/Temp/sisae2.xlsx', data_only=True)
ws = wb['송파구']

data = {}
for row in range(6, ws.max_row + 1):
    code = ws.cell(row, 4).value
    if code in ALL_CODES:
        data.setdefault(code, []).append({
            'name':   ws.cell(row, 6).value,
            'pyeong': ws.cell(row, 7).value,
            'mae':    ws.cell(row, 8).value,
            'jun':    ws.cell(row, 9).value,
            'sedae':  ws.cell(row, 13).value,
            'year':   ws.cell(row, 12).value,
        })

# ── PPT 로드 ───────────────────────────────────────────────────
PPT_PATH = 'C:/Users/한나/OneDrive/강성업무용/바탕 화면/문서/CLAUDE/송파구(입지분석_야자수그늘)_시세지도.pptx'
prs = Presentation(PPT_PATH)

HDR_H    = int(Pt(14))
ROW_H    = int(Pt(13))
ROW_H_EX = int(Pt(11))
T_INS    = 12000
FONT     = Pt(10)
FONT_EX  = Pt(9)

def add_boxes(slide, code, cx, box_top, rows):
    r0    = rows[0]
    name  = r0['name'] or f'코드{code}'
    year  = yr2(r0['year'])
    sedae = sd_str(r0['sedae'])
    hdr_text = f'{name} {year} {sedae}^'

    price_lines = []
    for r in rows:
        mae_raw = r['mae'] if not isinstance(r['mae'], str) else None
        jun_raw = r['jun'] if not isinstance(r['jun'], str) else None
        mae = fix(mae_raw)
        jun = fix(jun_raw)
        if mae is None: continue
        txt = (f'{r["pyeong"]}p {num(mae)}/{num(jun)} ({num(mae-jun)})'
               if jun else f'{r["pyeong"]}p {num(mae)}/-')
        bg, fg = get_color(mae)
        price_lines.append((txt, bg, fg))

    if not price_lines:
        price_lines = [('가격 미확인', '595959', 'FFFFFF')]

    extra_lines = get_extra_lines(name)

    all_texts = [hdr_text] + [t for t, _, _ in price_lines] + extra_lines
    BOX_W = calc_box_w(all_texts, font_pt=10, min_pt=70, padding_pt=14)

    L = cx - BOX_W // 2
    if L < int(Pt(3)): L = int(Pt(3))
    if L + BOX_W > SLIDE_W - int(Pt(3)):
        L = SLIDE_W - BOX_W - int(Pt(3))

    n_price = len(price_lines)
    n_extra = len(extra_lines)
    content_h = n_price * ROW_H + n_extra * ROW_H_EX + T_INS * 2
    total_h   = HDR_H + content_h

    T = box_top
    if T < int(Pt(3)): T = int(Pt(3))
    if T + total_h > SLIDE_H - int(Pt(3)):
        T = SLIDE_H - total_h - int(Pt(3))

    hdr_box = slide.shapes.add_textbox(L, T, BOX_W, HDR_H)
    hdr_box.fill.solid()
    hdr_box.fill.fore_color.rgb = RGBColor.from_string('111111')
    tf = hdr_box.text_frame
    tf.word_wrap = False
    set_bPr(tf, l=36000, r=36000, t=0, b=0, anchor='ctr')
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = hdr_text
    run.font.size = FONT
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string('FFFFFF')

    price_box = slide.shapes.add_textbox(L, T + HDR_H, BOX_W, content_h)
    set_no_fill(price_box)
    set_border(price_box, '333333', 1.0)
    tf2 = price_box.text_frame
    tf2.word_wrap = False
    set_bPr(tf2, l=36000, r=36000, t=T_INS, b=T_INS, anchor='t')

    for i, (txt, bg, fg) in enumerate(price_lines):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run2 = para.add_run()
        run2.text = txt
        run2.font.size = FONT
        run2.font.bold = True
        set_highlight(run2, bg, fg)

    for extra in extra_lines:
        para = tf2.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run3 = para.add_run()
        run3.text = extra
        run3.font.size = FONT_EX
        run3.font.bold = False
        set_highlight(run3, 'FFFFFF', '333333')

    return name, year, n_price, n_extra, BOX_W

def process_slide(slide, codes, slide_num):
    print(f'\n=== 슬라이드 {slide_num} 처리 ===')

    # 코드 박스 위치 수집 (중복 코드 포함 → 리스트로)
    code_positions = {}   # code → [(left, top, width, height), ...]
    delete_elems   = []

    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        text = sh.text_frame.text.strip()
        try:
            val = int(text)
            if val in codes:
                code_positions.setdefault(val, []).append(
                    (sh.left, sh.top, sh.width, sh.height))
                delete_elems.append(sh._element)
        except:
            pass

    found = sum(len(v) for v in code_positions.values())
    print(f'코드 박스 {found}개 발견 → 삭제')
    for el in delete_elems:
        parent = el.getparent()
        if parent is not None: parent.remove(el)

    # 박스 추가
    for code in codes:
        positions = code_positions.get(code, [])
        if not positions:
            print(f'  [{code}] 슬라이드에 없음')
            continue
        rows = data.get(code, [])
        if not rows:
            print(f'  [{code}] Excel 없음')
            continue

        rows_valid = [r for r in rows if fix(r['mae'] if not isinstance(r['mae'], str) else None) is not None]
        n_est  = max(len(rows_valid), 1)
        ex_est = len(get_extra_lines(rows[0]['name'] or ''))
        total_h_est = HDR_H + n_est * ROW_H + ex_est * ROW_H_EX + T_INS * 2

        # 중복 코드 → 각 위치마다 박스
        for sl, st, sw, sh_h in positions:
            cx = sl + sw // 2
            T = st - total_h_est - int(Pt(5))
            if T < int(Pt(3)):
                T = st + sh_h + int(Pt(5))
            name, year, np, ne, bw = add_boxes(slide, code, cx, T, rows)
            print(f'  [{code}] {name} {year}: {np}가격+{ne}추가 / 폭={bw//12700:.0f}pt')

# ── 슬라이드 16 (index 15) ─────────────────────────────────────
process_slide(prs.slides[15], S16_CODES, 16)

# ── 슬라이드 17 (index 16) ─────────────────────────────────────
process_slide(prs.slides[16], S17_CODES, 17)

# ── 저장 ────────────────────────────────────────────────────────
prs.save(PPT_PATH)
print(f'\n저장 완료: {PPT_PATH}')
shutil.copy2(PPT_PATH, 'C:/Temp/songpa_work.pptx')
print('복사 완료: C:/Temp/songpa_work.pptx')
