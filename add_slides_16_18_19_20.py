
"""
슬라이드 16(코드611) / 18 / 19 / 20 시세박스 + 연식표 채우기
"""
import re, json, shutil, openpyxl
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

A    = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
SLIDE_W = 9144000; SLIDE_H = 6858000

RAINBOW = {
    35:('000000','FFFFFF'), 34:('1A1A1A','FFFFFF'), 33:('2B2B2B','FFFFFF'),
    32:('424242','FFFFFF'), 31:('616161','FFFFFF'), 30:('757575','FFFFFF'),
    29:('4A0000','FFFFFF'), 28:('6A0F0F','FFFFFF'), 27:('8D1515','FFFFFF'),
    26:('C62828','FFFFFF'), 25:('E53935','FFFFFF'), 24:('F4511E','FFFFFF'),
    23:('FB8C00','FFFFFF'), 22:('FFB300','FFFFFF'), 21:('F9E400','000000'),
    20:('9CCC65','000000'), 19:('43A047','FFFFFF'), 18:('00897B','FFFFFF'),
    17:('00ACC1','000000'), 16:('29B6F6','000000'), 15:('1E88E5','FFFFFF'),
    14:('1565C0','FFFFFF'), 13:('1A237E','FFFFFF'), 12:('7B2FBE','FFFFFF'),
    11:('B07FE0','FFFFFF'), 10:('D7B8F5','000000'),
}

def fix(v):
    if v is None or isinstance(v, str): return None
    return v // 10 if v > 1000000 else v

def num(m):
    if m is None: return '-'
    v = m / 10000; s = f'{v:.1f}'
    return s[:-2] if s.endswith('.0') else s

def get_color(mae):
    if mae is None: return ('595959','FFFFFF')
    eok = int(mae // 10000)
    if eok >= 35: return RAINBOW[35]
    return RAINBOW.get(eok, ('595959','FFFFFF'))

def yr2(y):
    if not y: return '??'
    yr = re.sub(r'[^\d].*', '', str(y))
    return yr[-2:] if len(yr) >= 2 else yr

def sd_str(s):
    if s is None: return '?'
    return str(int(s)) if isinstance(s, float) else str(s)

def text_width_pt(text, font_pt=10):
    w = 0
    for c in text:
        if '가' <= c <= '힣': w += font_pt * 1.0
        elif c == ' ':        w += font_pt * 0.3
        elif c in '()':       w += font_pt * 0.4
        elif c in '/.,':      w += font_pt * 0.35
        else:                 w += font_pt * 0.6
    return w

def calc_box_w(lines, font_pt=10, min_pt=70, padding_pt=14):
    return max(int(Pt(max(text_width_pt(l, font_pt) for l in lines)+padding_pt)), int(Pt(min_pt)))

def set_highlight(run, bg_hex, fg_hex):
    rPr = run._r.find(f'{{{A}}}rPr')
    if rPr is None: rPr = etree.SubElement(run._r, f'{{{A}}}rPr')
    for tag in [f'{{{A}}}highlight', f'{{{A}}}solidFill']:
        for el in rPr.findall(tag): rPr.remove(el)
    sf = etree.SubElement(rPr, f'{{{A}}}solidFill')
    etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', fg_hex)
    hl = etree.SubElement(rPr, f'{{{A}}}highlight')
    etree.SubElement(hl, f'{{{A}}}srgbClr').set('val', bg_hex)

def set_no_fill(shape):
    sp = shape._element
    spPr = sp.find(f'{{{P_NS}}}spPr') or sp.find(f'{{{A}}}spPr')
    if spPr is None: return
    for tag in [f'{{{A}}}solidFill', f'{{{A}}}noFill', f'{{{A}}}gradFill']:
        for el in spPr.findall(tag): spPr.remove(el)
    etree.SubElement(spPr, f'{{{A}}}noFill')

def set_border(shape, hex_color='333333', width_pt=1.0):
    sp = shape._element
    spPr = sp.find(f'{{{P_NS}}}spPr') or sp.find(f'{{{A}}}spPr')
    if spPr is None: return
    for old in spPr.findall(f'{{{A}}}ln'): spPr.remove(old)
    ln = etree.SubElement(spPr, f'{{{A}}}ln')
    ln.set('w', str(int(Pt(width_pt))))
    sf = etree.SubElement(ln, f'{{{A}}}solidFill')
    etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', hex_color)

def set_bPr(tf, l=36000, r=36000, t=18000, b=18000, anchor='t'):
    bPr = tf._txBody.find(f'{{{A}}}bodyPr')
    if bPr is not None:
        bPr.set('lIns',str(l)); bPr.set('rIns',str(r))
        bPr.set('tIns',str(t)); bPr.set('bIns',str(b))
        bPr.set('anchor', anchor)

# ── hanpan5 extra ──────────────────────────────────────────────
def load_extra_info(html_path):
    with open(html_path, 'rb') as f:
        content = f.read().decode('utf-8')
    m = re.search(r'apartments=(\[.*?\]);', content, re.DOTALL)
    if not m: return {}
    arr = json.loads(m.group(1))
    price_pat = re.compile(r'^\d+p\(')
    result = {}
    for apt in arr:
        name_raw = apt.get('name','').split('\n')[0].strip()
        lines = apt.get('text','').split('\n')
        extra = [l.strip() for l in lines[1:] if l.strip() and not price_pat.match(l.strip())]
        filtered = [l for l in extra if not re.match(r"^\d{2}'\s*[\d,]+\^", l)]
        result[name_raw] = filtered
    return result

extra_db = load_extra_info('C:/Temp/hanpan5.html')

def get_extra_lines(apt_name):
    best, best_len = None, 0
    for html_name, lines in extra_db.items():
        for n in range(min(10, len(apt_name)), 3, -1):
            if apt_name[:n] in html_name and n > best_len:
                best_len = n; best = lines; break
    return best or []

# ── Excel 로드 ──────────────────────────────────────────────────
CODES_16  = [611]
CODES_18  = [128379,128378,139920,132624,139917,104122,107683,104126,
             25929,25930,26479,26222,26472,26481,26106,26471,26107,
             105367,107339,107550,109174,107807]
CODES_19  = [352,353,951,952,127132,355,102313,102314,119875]
CODES_20  = [648,643,646,416,115046,649,650,422,625]
ALL_CODES = list(set(CODES_16+CODES_18+CODES_19+CODES_20))

wb = openpyxl.load_workbook('C:/Temp/sisae2.xlsx', data_only=True)
ws = wb['송파구']
data = {}
for row in range(6, ws.max_row+1):
    code = ws.cell(row, 4).value
    if code in ALL_CODES:
        data.setdefault(code, []).append({
            'name':   ws.cell(row, 6).value,
            'pyeong': ws.cell(row, 7).value,
            'mae':    ws.cell(row, 8).value,
            'jun':    ws.cell(row, 9).value,
            'sedae':  ws.cell(row, 13).value,
            'year':   ws.cell(row, 12).value,
        })

PPT_PATH = 'C:/Users/한나/OneDrive/강성업무용/바탕 화면/문서/CLAUDE/송파구(입지분석_야자수그늘)_시세지도.pptx'
prs = Presentation(PPT_PATH)

HDR_H=int(Pt(14)); ROW_H=int(Pt(13)); ROW_H_EX=int(Pt(11)); T_INS=12000
FONT=Pt(10); FONT_EX=Pt(9)

# ── 박스 생성 함수 ─────────────────────────────────────────────
def add_boxes(slide, code, cx, box_top, rows):
    r0    = rows[0]
    name  = r0['name'] or f'코드{code}'
    year  = yr2(r0['year'])
    sedae = sd_str(r0['sedae'])
    hdr_text = f'{name} {year} {sedae}^'

    price_lines = []
    for r in rows:
        mae = fix(r['mae'] if not isinstance(r['mae'],str) else None)
        jun = fix(r['jun'] if not isinstance(r['jun'],str) else None)
        if mae is None: continue
        txt = (f'{r["pyeong"]}p {num(mae)}/{num(jun)} ({num(mae-jun)})'
               if jun else f'{r["pyeong"]}p {num(mae)}/-')
        bg, fg = get_color(mae)
        price_lines.append((txt, bg, fg))

    if not price_lines:
        price_lines = [('가격 미확인', '595959', 'FFFFFF')]

    extra_lines = get_extra_lines(name)
    all_texts = [hdr_text]+[t for t,_,_ in price_lines]+extra_lines
    BOX_W = calc_box_w(all_texts, font_pt=10, min_pt=70, padding_pt=14)

    L = max(int(Pt(3)), cx - BOX_W//2)
    if L + BOX_W > SLIDE_W - int(Pt(3)): L = SLIDE_W - BOX_W - int(Pt(3))

    n_price = len(price_lines); n_extra = len(extra_lines)
    content_h = n_price*ROW_H + n_extra*ROW_H_EX + T_INS*2
    total_h = HDR_H + content_h

    T = box_top
    if T < int(Pt(3)): T = int(Pt(3))
    if T + total_h > SLIDE_H - int(Pt(3)): T = SLIDE_H - total_h - int(Pt(3))

    hdr_box = slide.shapes.add_textbox(L, T, BOX_W, HDR_H)
    hdr_box.fill.solid(); hdr_box.fill.fore_color.rgb = RGBColor.from_string('111111')
    tf = hdr_box.text_frame; tf.word_wrap = False
    set_bPr(tf, l=36000, r=36000, t=0, b=0, anchor='ctr')
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = hdr_text
    run.font.size = FONT; run.font.bold = True
    run.font.color.rgb = RGBColor.from_string('FFFFFF')

    price_box = slide.shapes.add_textbox(L, T+HDR_H, BOX_W, content_h)
    set_no_fill(price_box); set_border(price_box,'333333',1.0)
    tf2 = price_box.text_frame; tf2.word_wrap = False
    set_bPr(tf2, l=36000, r=36000, t=T_INS, b=T_INS, anchor='t')

    for i,(txt,bg,fg) in enumerate(price_lines):
        para = tf2.paragraphs[0] if i==0 else tf2.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run2 = para.add_run(); run2.text = txt
        run2.font.size = FONT; run2.font.bold = True
        set_highlight(run2, bg, fg)

    for extra in extra_lines:
        para = tf2.add_paragraph(); para.alignment = PP_ALIGN.CENTER
        run3 = para.add_run(); run3.text = extra
        run3.font.size = FONT_EX; run3.font.bold = False
        set_highlight(run3, 'FFFFFF', '333333')

    return name, year, n_price, n_extra, BOX_W

# ── 슬라이드별 처리 함수 ──────────────────────────────────────
def process_slide(slide, codes, slide_num, table_idx):
    print(f'\n=== 슬라이드 {slide_num} ===')

    # 코드박스 위치 수집 (텍스트박스 + 단일셀 표)
    code_positions = {}  # code → [(l,t,w,h), ...]
    delete_elems   = []

    for sh in slide.shapes:
        # 텍스트박스 코드
        if sh.has_text_frame:
            try:
                val = int(sh.text_frame.text.strip())
                if val in codes:
                    code_positions.setdefault(val,[]).append(
                        (sh.left, sh.top, sh.width, sh.height))
                    delete_elems.append(sh._element)
            except: pass
        # 단일셀 표 코드 (1행×1열)
        elif sh.has_table:
            tbl = sh.table
            if len(tbl.rows)==1 and len(tbl.columns)==1:
                try:
                    val = int(tbl.cell(0,0).text.strip())
                    if val in codes:
                        code_positions.setdefault(val,[]).append(
                            (sh.left, sh.top, sh.width, sh.height))
                        delete_elems.append(sh._element)
                except: pass

    print(f'코드박스 {sum(len(v) for v in code_positions.values())}개 → 삭제')
    for el in delete_elems:
        parent = el.getparent()
        if parent is not None: parent.remove(el)

    # 박스 추가
    for code in codes:
        positions = code_positions.get(code, [])
        if not positions:
            print(f'  [{code}] 슬라이드에 없음'); continue
        rows = data.get(code, [])
        if not rows:
            print(f'  [{code}] Excel 없음'); continue

        rows_valid = [r for r in rows if fix(r['mae'] if not isinstance(r['mae'],str) else None) is not None]
        n_est  = max(len(rows_valid), 1)
        ex_est = len(get_extra_lines(rows[0]['name'] or ''))
        total_h_est = HDR_H + n_est*ROW_H + ex_est*ROW_H_EX + T_INS*2

        for sl, st, sw, sh_h in positions:
            cx = sl + sw//2
            T = st - total_h_est - int(Pt(5))
            if T < int(Pt(3)): T = st + sh_h + int(Pt(5))
            name, year, np, ne, bw = add_boxes(slide, code, cx, T, rows)
            print(f'  [{code}] {name} {year}: {np}가격+{ne}추가 / 폭={bw//12700:.0f}pt')

    # 연식 표 채우기
    fill_era_table(slide, table_idx, slide_num)

# ── 연식 표 채우기 ────────────────────────────────────────────
HDR_PAT   = re.compile(r'^(.+?)\s+(\d{2})\s+[\d,]+\^', re.DOTALL)
PRICE_PAT = re.compile(r'^(\d+)p\s+([\d.]+)/([\d.]+|-)')
ERA_RANGES = [(2019,2099),(2014,2018),(2004,2013),(0,2003)]
FONT_SZ = Pt(10)

def parse_header(text):
    text2 = text.replace('\n',' ').strip()
    m = HDR_PAT.match(text2)
    if not m: return None, None
    yr = int(m.group(2))
    return m.group(1).strip(), (2000+yr if yr<=30 else 1900+yr)

def get_best_price(txt, target=84):
    best_diff, best_str = float('inf'), None
    for line in txt.split('\n'):
        m = PRICE_PAT.match(line.strip())
        if m:
            p = int(m.group(1)); diff = abs(p-target)
            if diff < best_diff:
                best_diff = diff
                best_str = f'{p}p {m.group(2)}/{m.group(3)}'
    return best_str

def get_txBody(cell):
    return cell._tc.find(f'{{{A}}}txBody')

def clear_cell(cell):
    txBody = get_txBody(cell)
    if txBody is None: return
    for p in list(txBody.findall(f'{{{A}}}p')): txBody.remove(p)

def add_para(cell, text, color='111111'):
    txBody = get_txBody(cell)
    if txBody is None: return
    p_el = etree.SubElement(txBody, f'{{{A}}}p')
    pPr  = etree.SubElement(p_el, f'{{{A}}}pPr')
    pPr.set('algn','l')
    lnSpc = etree.SubElement(pPr, f'{{{A}}}lnSpc')
    etree.SubElement(lnSpc, f'{{{A}}}spcPct').set('val','100000')
    for tag in ['spcBef','spcAft']:
        el = etree.SubElement(pPr, f'{{{A}}}{tag}')
        etree.SubElement(el, f'{{{A}}}spcPts').set('val','0')
    r_el = etree.SubElement(p_el, f'{{{A}}}r')
    rPr  = etree.SubElement(r_el, f'{{{A}}}rPr')
    rPr.set('lang','ko-KR'); rPr.set('sz',str(int(FONT_SZ.pt*100)))
    rPr.set('b','0'); rPr.set('dirty','0')
    sf = etree.SubElement(rPr, f'{{{A}}}solidFill')
    etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', color)
    etree.SubElement(r_el, f'{{{A}}}t').text = text

def sort_key(apt):
    m = re.search(r'p\s+([\d.]+)/', apt[2])
    return (apt[1], -(float(m.group(1)) if m else 0.0))

def find_era_table(slide):
    """5행×3열 연식표 찾기 (동적)"""
    for sh in slide.shapes:
        if sh.has_table:
            tbl = sh.table
            if len(tbl.rows)==5 and len(tbl.columns)==3:
                h0 = tbl.cell(0,0).text.strip()
                if '연식' in h0 or '단지' in tbl.cell(0,1).text.strip():
                    return tbl
    return None

def fill_era_table(slide, table_idx, slide_num):
    shapes = list(slide.shapes)
    header_shapes  = {}
    content_shapes = {}
    for sh in shapes:
        if not sh.has_text_frame: continue
        text = sh.text_frame.text.strip()
        if not text: continue
        name, year = parse_header(text)
        if name and year:
            header_shapes[id(sh)] = (sh, name, year)
        elif PRICE_PAT.search(text) or '가격미확인' in text or '가격 미확인' in text:
            content_shapes[id(sh)] = (sh, text)

    print(f'  헤더 {len(header_shapes)}개, 가격 {len(content_shapes)}개')

    apartments = []
    seen = set()
    for cid, (csh, ctext) in content_shapes.items():
        best_hsh, best_dist = None, float('inf')
        for hid, (hsh, hname, hyear) in header_shapes.items():
            dy = csh.top - hsh.top
            dx = abs(csh.left - hsh.left)
            if dy > 0 and dx < int(Pt(120)):
                dist = dy + dx*0.3
                if dist < best_dist:
                    best_dist = dist; best_hsh = (hsh, hname, hyear)
        if best_hsh:
            _, name, year = best_hsh
            key = (name, year)
            if key in seen: continue
            seen.add(key)
            price_str = get_best_price(ctext) or '가격미확인'
            apartments.append((name, year, price_str))

    era_apts = [[] for _ in range(4)]
    for name, year, price_str in apartments:
        for i, (y0,y1) in enumerate(ERA_RANGES):
            if y0 <= year <= y1:
                era_apts[i].append((name, year, price_str)); break

    for i in range(4): era_apts[i].sort(key=sort_key)

    # 표 채우기 (동적 탐색)
    tbl = find_era_table(slide)
    if tbl is None:
        print(f'  [오류] 연식표 못 찾음'); return
    for row_idx, apts in enumerate(era_apts):
        row_num = row_idx + 1
        dcell = tbl.cell(row_num, 1); pcell = tbl.cell(row_num, 2)
        clear_cell(dcell); clear_cell(pcell)
        if not apts:
            add_para(dcell, '-', color='888888')
            add_para(pcell, '-', color='888888'); continue
        for name, year, price in apts:
            add_para(dcell, f"{name}{str(year)[-2:]}'")
            add_para(pcell, price)
        print(f'  row{row_num}: {len(apts)}개')

# ── 실행 ──────────────────────────────────────────────────────
# 슬라이드 16: 코드611만 추가 + 표 재생성
process_slide(prs.slides[15], CODES_16,  16, table_idx=5)
# 슬라이드 18
process_slide(prs.slides[17], CODES_18,  18, table_idx=5)
# 슬라이드 19 (era table은 index 6이지만 shape 삭제 후 달라질 수 있어 이름으로 탐색)
# 먼저 실행 전 table_idx 동적 탐색
def find_era_table_idx(slide):
    for i, sh in enumerate(slide.shapes):
        if sh.has_table:
            tbl = sh.table
            if len(tbl.rows)==5 and len(tbl.columns)==3:
                if '연식' in tbl.cell(0,0).text or '단지' in tbl.cell(0,1).text:
                    return i
    return 6  # fallback

process_slide(prs.slides[18], CODES_19,  19, table_idx=find_era_table_idx(prs.slides[18]))
process_slide(prs.slides[19], CODES_20,  20, table_idx=find_era_table_idx(prs.slides[19]))

# 저장
prs.save(PPT_PATH)
shutil.copy2(PPT_PATH, 'C:/Temp/songpa_work.pptx')
print('\n저장 완료')
