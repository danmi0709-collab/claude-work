
"""코드 634 (슬라이드 16) 박스 추가 + 슬라이드 15 표 채우기"""
import re, json, shutil, openpyxl
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

A    = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
SLIDE_W = 9144000; SLIDE_H = 6858000

RAINBOW = {
    35:('000000','FFFFFF'), 34:('1A1A1A','FFFFFF'), 33:('2B2B2B','FFFFFF'),
    32:('424242','FFFFFF'), 31:('616161','FFFFFF'), 30:('757575','FFFFFF'),
    29:('4A0000','FFFFFF'), 28:('6A0F0F','FFFFFF'), 27:('8D1515','FFFFFF'),
    26:('C62828','FFFFFF'), 25:('E53935','FFFFFF'), 24:('F4511E','FFFFFF'),
    23:('FB8C00','FFFFFF'), 22:('FFB300','FFFFFF'), 21:('F9E400','000000'),
    20:('9CCC65','000000'), 19:('43A047','FFFFFF'), 18:('00897B','FFFFFF'),
    17:('00ACC1','000000'), 16:('29B6F6','000000'), 15:('1E88E5','FFFFFF'),
    14:('1565C0','FFFFFF'), 13:('1A237E','FFFFFF'), 12:('7B2FBE','FFFFFF'),
    11:('B07FE0','FFFFFF'), 10:('D7B8F5','000000'),
}

def fix(v):
    if v is None or isinstance(v, str): return None
    return v // 10 if v > 1000000 else v

def num(m):
    if m is None: return '-'
    v = m / 10000; s = f'{v:.1f}'
    return s[:-2] if s.endswith('.0') else s

def get_color(mae):
    if mae is None: return ('595959','FFFFFF')
    eok = int(mae // 10000)
    if eok >= 35: return RAINBOW[35]
    return RAINBOW.get(eok, ('595959','FFFFFF'))

def yr2(y):
    if not y: return '??'
    yr = re.sub(r'[^\d].*', '', str(y))
    return yr[-2:] if len(yr) >= 2 else yr

def sd_str(s):
    if s is None: return '?'
    return str(int(s)) if isinstance(s, float) else str(s)

def text_width_pt(text, font_pt=10):
    w = 0
    for c in text:
        if '가' <= c <= '힣': w += font_pt * 1.0
        elif c == ' ':        w += font_pt * 0.3
        elif c in '()':       w += font_pt * 0.4
        elif c in '/.,':      w += font_pt * 0.35
        else:                 w += font_pt * 0.6
    return w

def calc_box_w(lines, font_pt=10, min_pt=70, padding_pt=14):
    return max(int(Pt(max(text_width_pt(l, font_pt) for l in lines) + padding_pt)), int(Pt(min_pt)))

def set_highlight(run, bg_hex, fg_hex):
    rPr = run._r.find(f'{{{A}}}rPr')
    if rPr is None: rPr = etree.SubElement(run._r, f'{{{A}}}rPr')
    for tag in [f'{{{A}}}highlight', f'{{{A}}}solidFill']:
        for el in rPr.findall(tag): rPr.remove(el)
    sf = etree.SubElement(rPr, f'{{{A}}}solidFill')
    etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', fg_hex)
    hl = etree.SubElement(rPr, f'{{{A}}}highlight')
    etree.SubElement(hl, f'{{{A}}}srgbClr').set('val', bg_hex)

def set_no_fill(shape):
    sp = shape._element
    spPr = sp.find(f'{{{P_NS}}}spPr')
    if spPr is None: spPr = sp.find(f'{{{A}}}spPr')
    if spPr is None: return
    for tag in [f'{{{A}}}solidFill', f'{{{A}}}noFill', f'{{{A}}}gradFill']:
        for el in spPr.findall(tag): spPr.remove(el)
    etree.SubElement(spPr, f'{{{A}}}noFill')

def set_border(shape, hex_color='333333', width_pt=1.0):
    sp = shape._element
    spPr = sp.find(f'{{{P_NS}}}spPr')
    if spPr is None: spPr = sp.find(f'{{{A}}}spPr')
    if spPr is None: return
    for old in spPr.findall(f'{{{A}}}ln'): spPr.remove(old)
    ln = etree.SubElement(spPr, f'{{{A}}}ln')
    ln.set('w', str(int(Pt(width_pt))))
    sf = etree.SubElement(ln, f'{{{A}}}solidFill')
    etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', hex_color)

def set_bPr(tf, l=36000, r=36000, t=18000, b=18000, anchor='t'):
    bPr = tf._txBody.find(f'{{{A}}}bodyPr')
    if bPr is not None:
        bPr.set('lIns',str(l)); bPr.set('rIns',str(r))
        bPr.set('tIns',str(t)); bPr.set('bIns',str(b))
        bPr.set('anchor', anchor)

# ── hanpan5 extra ──────────────────────────────────────────────
def load_extra_info(html_path):
    with open(html_path, 'rb') as f:
        content = f.read().decode('utf-8')
    m = re.search(r'apartments=(\[.*?\]);', content, re.DOTALL)
    if not m: return {}
    arr = json.loads(m.group(1))
    price_pat = re.compile(r'^\d+p\(')
    result = {}
    for apt in arr:
        name_raw = apt.get('name','').split('\n')[0].strip()
        lines = apt.get('text','').split('\n')
        extra = [l.strip() for l in lines[1:] if l.strip() and not price_pat.match(l.strip())]
        filtered = [l for l in extra if not re.match(r"^\d{2}'\s*[\d,]+\^", l)]
        result[name_raw] = filtered
    return result

extra_db = load_extra_info('C:/Temp/hanpan5.html')

def get_extra_lines(apt_name):
    best, best_len = None, 0
    for html_name, lines in extra_db.items():
        for n in range(min(10, len(apt_name)), 3, -1):
            if apt_name[:n] in html_name and n > best_len:
                best_len = n; best = lines; break
    return best or []

# ── Excel ──────────────────────────────────────────────────────
wb = openpyxl.load_workbook('C:/Temp/sisae2.xlsx', data_only=True)
ws = wb['송파구']
rows_634 = []
for row in range(6, ws.max_row + 1):
    if ws.cell(row, 4).value == 634:
        rows_634.append({
            'name':   ws.cell(row, 6).value,
            'pyeong': ws.cell(row, 7).value,
            'mae':    ws.cell(row, 8).value,
            'jun':    ws.cell(row, 9).value,
            'sedae':  ws.cell(row, 13).value,
            'year':   ws.cell(row, 12).value,
        })

PPT_PATH = 'C:/Users/한나/OneDrive/강성업무용/바탕 화면/문서/CLAUDE/송파구(입지분석_야자수그늘)_시세지도.pptx'
prs = Presentation(PPT_PATH)
slide16 = prs.slides[15]

HDR_H=int(Pt(14)); ROW_H=int(Pt(13)); ROW_H_EX=int(Pt(11)); T_INS=12000
FONT=Pt(10); FONT_EX=Pt(9)

# ── 코드 634 위치 & 삭제 ────────────────────────────────────────
code_pos = None; del_el = None
for sh in slide16.shapes:
    if not sh.has_text_frame: continue
    try:
        if int(sh.text_frame.text.strip()) == 634:
            code_pos = (sh.left, sh.top, sh.width, sh.height)
            del_el = sh._element
    except: pass

if del_el is not None:
    del_el.getparent().remove(del_el)
    print('코드 634 삭제')

# ── 박스 추가 ──────────────────────────────────────────────────
r0    = rows_634[0]
name  = r0['name']
year  = yr2(r0['year'])
sedae = sd_str(r0['sedae'])
hdr_text = f'{name} {year} {sedae}^'

price_lines = []
for r in rows_634:
    mae = fix(r['mae'] if not isinstance(r['mae'], str) else None)
    jun = fix(r['jun'] if not isinstance(r['jun'], str) else None)
    if mae is None: continue
    txt = (f'{r["pyeong"]}p {num(mae)}/{num(jun)} ({num(mae-jun)})'
           if jun else f'{r["pyeong"]}p {num(mae)}/-')
    bg, fg = get_color(mae)
    price_lines.append((txt, bg, fg))

extra_lines = get_extra_lines(name)
all_texts = [hdr_text] + [t for t,_,_ in price_lines] + extra_lines
BOX_W = calc_box_w(all_texts, font_pt=10, min_pt=70, padding_pt=14)

sl, st, sw, sh_h = code_pos
cx = sl + sw // 2
n_price = len(price_lines); n_extra = len(extra_lines)
content_h = n_price*ROW_H + n_extra*ROW_H_EX + T_INS*2
T = st - HDR_H - content_h - int(Pt(5))
if T < int(Pt(3)): T = st + sh_h + int(Pt(5))
L = max(int(Pt(3)), cx - BOX_W//2)
if L + BOX_W > SLIDE_W - int(Pt(3)): L = SLIDE_W - BOX_W - int(Pt(3))

hdr_box = slide16.shapes.add_textbox(L, T, BOX_W, HDR_H)
hdr_box.fill.solid(); hdr_box.fill.fore_color.rgb = RGBColor.from_string('111111')
tf = hdr_box.text_frame; tf.word_wrap = False
set_bPr(tf, l=36000, r=36000, t=0, b=0, anchor='ctr')
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = hdr_text
run.font.size = FONT; run.font.bold = True
run.font.color.rgb = RGBColor.from_string('FFFFFF')

price_box = slide16.shapes.add_textbox(L, T+HDR_H, BOX_W, content_h)
set_no_fill(price_box); set_border(price_box, '333333', 1.0)
tf2 = price_box.text_frame; tf2.word_wrap = False
set_bPr(tf2, l=36000, r=36000, t=T_INS, b=T_INS, anchor='t')

for i, (txt, bg, fg) in enumerate(price_lines):
    para = tf2.paragraphs[0] if i==0 else tf2.add_paragraph()
    para.alignment = PP_ALIGN.CENTER
    run2 = para.add_run(); run2.text = txt
    run2.font.size = FONT; run2.font.bold = True
    set_highlight(run2, bg, fg)

for extra in extra_lines:
    para = tf2.add_paragraph(); para.alignment = PP_ALIGN.CENTER
    run3 = para.add_run(); run3.text = extra
    run3.font.size = FONT_EX; run3.font.bold = False
    set_highlight(run3, 'FFFFFF', '333333')

print(f'[634] {name} {year}: {n_price}가격+{n_extra}추가 / 폭={BOX_W//12700:.0f}pt')

# ── 슬라이드 15 표 채우기 ──────────────────────────────────────
HDR_PAT   = re.compile(r'^(.+?)\s+(\d{2})\s+[\d,]+\^', re.DOTALL)
PRICE_PAT = re.compile(r'^(\d+)p\s+([\d.]+)/([\d.]+|-)')
ERA_RANGES = [(2019,2099),(2014,2018),(2004,2013),(0,2003)]
FONT_SZ = Pt(10)

def parse_header(text):
    text2 = text.replace('\n',' ').strip()
    m = HDR_PAT.match(text2)
    if not m: return None, None
    name = m.group(1).strip()
    yr   = int(m.group(2))
    return name, (2000+yr if yr<=30 else 1900+yr)

def get_best_price(content_text, target=84):
    best_diff, best_str = float('inf'), None
    for line in content_text.split('\n'):
        m = PRICE_PAT.match(line.strip())
        if m:
            p = int(m.group(1)); diff = abs(p-target)
            if diff < best_diff:
                best_diff = diff
                best_str = f'{p}p {m.group(2)}/{m.group(3)}'
    return best_str

def get_txBody(cell):
    return cell._tc.find(f'{{{A}}}txBody')

def clear_cell(cell):
    txBody = get_txBody(cell)
    if txBody is None: return
    for p in list(txBody.findall(f'{{{A}}}p')): txBody.remove(p)

def add_para(cell, text, color='111111'):
    txBody = get_txBody(cell)
    if txBody is None: return
    p_el = etree.SubElement(txBody, f'{{{A}}}p')
    pPr  = etree.SubElement(p_el, f'{{{A}}}pPr')
    pPr.set('algn','l')
    lnSpc = etree.SubElement(pPr, f'{{{A}}}lnSpc')
    etree.SubElement(lnSpc, f'{{{A}}}spcPct').set('val','100000')
    spcBef = etree.SubElement(pPr, f'{{{A}}}spcBef')
    etree.SubElement(spcBef, f'{{{A}}}spcPts').set('val','0')
    spcAft = etree.SubElement(pPr, f'{{{A}}}spcAft')
    etree.SubElement(spcAft, f'{{{A}}}spcPts').set('val','0')
    r_el = etree.SubElement(p_el, f'{{{A}}}r')
    rPr  = etree.SubElement(r_el, f'{{{A}}}rPr')
    rPr.set('lang','ko-KR'); rPr.set('sz', str(int(FONT_SZ.pt*100)))
    rPr.set('b','0'); rPr.set('dirty','0')
    sf = etree.SubElement(rPr, f'{{{A}}}solidFill')
    etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', color)
    etree.SubElement(r_el, f'{{{A}}}t').text = text

slide15 = prs.slides[14]
shapes15 = list(slide15.shapes)

header_shapes  = {}
content_shapes = {}
for sh in shapes15:
    if not sh.has_text_frame: continue
    text = sh.text_frame.text.strip()
    if not text: continue
    name, year = parse_header(text)
    if name and year:
        header_shapes[id(sh)] = (sh, name, year)
    elif PRICE_PAT.search(text) or '가격미확인' in text or '가격 미확인' in text:
        content_shapes[id(sh)] = (sh, text)

print(f'\n슬라이드15 헤더 {len(header_shapes)}개, 가격 {len(content_shapes)}개')

apartments = []
seen = set()
for cid, (csh, ctext) in content_shapes.items():
    best_hsh, best_dist = None, float('inf')
    for hid, (hsh, hname, hyear) in header_shapes.items():
        dy = csh.top - hsh.top
        dx = abs(csh.left - hsh.left)
        if dy > 0 and dx < int(Pt(120)):
            dist = dy + dx*0.3
            if dist < best_dist:
                best_dist = dist; best_hsh = (hsh, hname, hyear)
    if best_hsh:
        _, name, year = best_hsh
        key = (name, year)
        if key in seen: continue   # 중복 제거
        seen.add(key)
        price_str = get_best_price(ctext) or '가격미확인'
        apartments.append((name, year, price_str))

era_apts = [[] for _ in range(4)]
for name, year, price_str in apartments:
    for i, (y0,y1) in enumerate(ERA_RANGES):
        if y0 <= year <= y1:
            era_apts[i].append((name, year, price_str)); break

def sort_key(apt):
    m = re.search(r'p\s+([\d.]+)/', apt[2])
    return (apt[1], -(float(m.group(1)) if m else 0.0))

for i in range(4): era_apts[i].sort(key=sort_key)

print('정렬 결과:')
for i, apts in enumerate(era_apts):
    for name, year, price in apts:
        print(f'  row{i+1} | {name}{str(year)[-2:]}\' | {price}')

tbl = shapes15[5].table
for row_idx, apts in enumerate(era_apts):
    row_num = row_idx + 1
    dcell = tbl.cell(row_num, 1); pcell = tbl.cell(row_num, 2)
    clear_cell(dcell); clear_cell(pcell)
    if not apts:
        add_para(dcell, '-', color='888888'); add_para(pcell, '-', color='888888'); continue
    for name, year, price in apts:
        add_para(dcell, f"{name}{str(year)[-2:]}'")
        add_para(pcell, price)
    print(f'  표 row{row_num} 채움: {len(apts)}개')

# 저장
prs.save(PPT_PATH)
shutil.copy2(PPT_PATH, 'C:/Temp/songpa_work.pptx')
print('\n저장 완료')
