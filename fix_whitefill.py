
"""
모든 처리된 슬라이드(14~20)의 가격 컨텐츠 박스에 흰색 배경 채우기
- 헤더박스(검정)는 그대로
- 가격/추가정보 박스만 흰색으로
"""
import re, shutil
from pptx import Presentation
from lxml import etree

A    = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

PRICE_PAT = re.compile(r'^\d+p\s+[\d.]+/', re.MULTILINE)

def set_white_fill(shape):
    sp = shape._element
    spPr = sp.find(f'{{{P_NS}}}spPr')
    if spPr is None: spPr = sp.find(f'{{{A}}}spPr')
    if spPr is None: return False
    for tag in [f'{{{A}}}solidFill', f'{{{A}}}noFill', f'{{{A}}}gradFill']:
        for el in spPr.findall(tag): spPr.remove(el)
    sf = etree.SubElement(spPr, f'{{{A}}}solidFill')
    etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', 'FFFFFF')
    return True

PPT_PATH = 'C:/Users/한나/OneDrive/강성업무용/바탕 화면/문서/CLAUDE/송파구(입지분석_야자수그늘)_시세지도.pptx'
prs = Presentation(PPT_PATH)

# 슬라이드 14~20 (index 13~19)
SLIDE_INDICES = list(range(13, 20))

total = 0
for idx in SLIDE_INDICES:
    slide = prs.slides[idx]
    count = 0
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        text = sh.text_frame.text.strip()
        if not text: continue
        # 가격박스: PRICE_PAT 또는 가격미확인 포함 AND 헤더형식(^) 아님
        is_price_box = (PRICE_PAT.search(text) or '가격미확인' in text or '가격 미확인' in text)
        is_header    = text.endswith('^')
        if is_price_box and not is_header:
            if set_white_fill(sh):
                count += 1
    print(f'슬라이드 {idx+1}: {count}개 박스 흰색 채우기')
    total += count

prs.save(PPT_PATH)
shutil.copy2(PPT_PATH, 'C:/Temp/songpa_work.pptx')
print(f'\n총 {total}개 / 저장 완료')
