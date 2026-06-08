import sys, io, json
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import math

# 이미지 크기
img_path = '송파구/송파구이미지2.png'
img = Image.open(img_path)
IMG_W, IMG_H = img.size
print(f"이미지: {IMG_W} x {IMG_H}")

# 슬라이드 크기 - 이미지 비율 맞춰서 큼지막하게
# 가로 13.5인치 기준, 세로는 비율 계산
SLIDE_W_IN = 13.5
SLIDE_H_IN = SLIDE_W_IN * (IMG_H / IMG_W)
print(f"슬라이드: {SLIDE_W_IN:.2f} x {SLIDE_H_IN:.2f} 인치")

prs = Presentation()
prs.slide_width = Inches(SLIDE_W_IN)
prs.slide_height = Inches(SLIDE_H_IN)

# 빈 슬라이드 추가
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# 이미지 배경으로 깔기 (슬라이드 전체)
slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

# 데이터 로드
with open('송파구/apt_data.json', encoding='utf-8') as f:
    apts = json.load(f)
print(f"아파트: {len(apts)}개")

# hex → RGBColor
def hex_to_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# 좌표 → 슬라이드 위치
LAT_MAX, LAT_MIN = 37.5470, 37.4570
LNG_MIN, LNG_MAX = 127.0750, 127.1970

def to_slide_pos(lat, lng):
    rx = (lng - LNG_MIN) / (LNG_MAX - LNG_MIN)
    ry = (LAT_MAX - lat) / (LAT_MAX - LAT_MIN)
    return rx * SLIDE_W_IN, ry * SLIDE_H_IN

# 텍스트박스 크기
BOX_W_IN = 1.1
BOX_H_IN = 0.45

for apt in apts:
    sx, sy = to_slide_pos(apt['lat'], apt['lng'])
    left = Inches(sx - BOX_W_IN/2)
    top = Inches(sy - BOX_H_IN/2)
    width = Inches(BOX_W_IN)
    height = Inches(BOX_H_IN)

    # 둥근 사각형 도형 추가
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

    # 채우기 색
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(apt['bg'])

    # 테두리
    line = shape.line
    line.color.rgb = hex_to_rgb(apt['color'])
    line.width = Pt(1.5)

    # 텍스트
    tf = shape.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 1줄: 아파트명
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = apt['name']
    r1.font.name = '맑은 고딕'
    r1.font.size = Pt(8)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0x11, 0x11, 0x11)

    # 2줄: 연식·세대수
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"{apt['yr2']}' {apt['units']:,}^"
    r2.font.name = '맑은 고딕'
    r2.font.size = Pt(7)
    r2.font.bold = False
    r2.font.color.rgb = hex_to_rgb(apt['color'])

# 범례 추가 (왼쪽 상단)
legend_data = [
    ('~1989 구축',       '#c0392b', '#fde8e8'),
    ('1990~1999',        '#d35400', '#fef0e0'),
    ('2000~2009',        '#1e8449', '#e8f8ee'),
    ('2010~2019',        '#1a5276', '#e8f2fd'),
    ('2020년~',          '#6c3483', '#f2e8fd'),
]
LEG_X, LEG_Y = 0.15, 0.15
LEG_W, LEG_H = 1.4, 0.3
for i, (label, color, bg) in enumerate(legend_data):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(LEG_X), Inches(LEG_Y + i*(LEG_H + 0.05)),
        Inches(LEG_W), Inches(LEG_H)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = hex_to_rgb(bg)
    sh.line.color.rgb = hex_to_rgb(color)
    sh.line.width = Pt(1.5)
    tf = sh.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = '맑은 고딕'
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = hex_to_rgb(color)

# 안내 텍스트
note = slide.shapes.add_textbox(Inches(0.15), Inches(LEG_Y + 5*(LEG_H+0.05) + 0.1),
                                 Inches(2.0), Inches(0.5))
ntf = note.text_frame
ntf.word_wrap = True
np_ = ntf.paragraphs[0]
nr = np_.add_run()
nr.text = "※ 라벨을 드래그해서\n   원하는 위치로 옮기세요"
nr.font.name = '맑은 고딕'
nr.font.size = Pt(9)
nr.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# ===== 두번째 슬라이드: 가나다순 라벨 목록 =====
slide2 = prs.slides.add_slide(blank_layout)

# 제목
title_box = slide2.shapes.add_textbox(Inches(0.3), Inches(0.15),
                                       Inches(SLIDE_W_IN - 0.6), Inches(0.5))
tt = title_box.text_frame
tp = tt.paragraphs[0]
tr = tp.add_run()
tr.text = "송파구 아파트 라벨 (가나다순) - 필요한 라벨을 1번 슬라이드로 복사·드래그하세요"
tr.font.name = '맑은 고딕'
tr.font.size = Pt(14)
tr.font.bold = True
tr.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

# 가나다순 정렬
apts_sorted = sorted(apts, key=lambda x: x['name'])

# 그리드 배치: 8열 × 15행
COLS = 8
START_X = 0.25
START_Y = 0.75
GAP_X = 0.05
GAP_Y = 0.10
CELL_W = (SLIDE_W_IN - START_X*2 - GAP_X*(COLS-1)) / COLS
CELL_H = 0.55

for i, apt in enumerate(apts_sorted):
    col = i % COLS
    row = i // COLS
    left = Inches(START_X + col * (CELL_W + GAP_X))
    top = Inches(START_Y + row * (CELL_H + GAP_Y))
    width = Inches(CELL_W)
    height = Inches(CELL_H)

    shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(apt['bg'])
    shape.line.color.rgb = hex_to_rgb(apt['color'])
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = apt['name']
    r1.font.name = '맑은 고딕'
    r1.font.size = Pt(8)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0x11, 0x11, 0x11)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"{apt['yr2']}' {apt['units']:,}^"
    r2.font.name = '맑은 고딕'
    r2.font.size = Pt(7)
    r2.font.color.rgb = hex_to_rgb(apt['color'])

# ===== 세번째 슬라이드: 송파구 지하철역 (호선별 색상) =====
slide3 = prs.slides.add_slide(blank_layout)

# 제목
title_box = slide3.shapes.add_textbox(Inches(0.3), Inches(0.2),
                                       Inches(SLIDE_W_IN - 0.6), Inches(0.5))
tt = title_box.text_frame
tp = tt.paragraphs[0]
tr = tp.add_run()
tr.text = "송파구 지하철역 (호선별)"
tr.font.name = '맑은 고딕'
tr.font.size = Pt(20)
tr.font.bold = True
tr.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

# 호선별 데이터 (서울교통공사 공식 색상)
lines = [
    {'name': '2호선', 'color': '#00A84D', 'stations': ['잠실나루', '잠실', '잠실새내']},
    {'name': '3호선', 'color': '#EF7C1C', 'stations': ['가락시장', '경찰병원', '오금']},
    {'name': '5호선', 'color': '#996CAC', 'stations': ['방이', '오금', '개롱', '거여', '마천']},
    {'name': '8호선', 'color': '#E6186C', 'stations': ['잠실', '몽촌토성', '석촌', '송파', '가락시장', '문정', '장지']},
    {'name': '9호선', 'color': '#BDB092', 'stations': ['석촌고분', '송파나루', '한성백제', '올림픽공원']},
]

ROW_START_Y = 1.0
ROW_H = 1.2          # 한 호선이 차지하는 세로 공간
LINE_LABEL_W = 1.2   # 호선 이름 박스 너비
STATION_W = 1.3      # 역 박스 너비
STATION_H = 0.6      # 역 박스 높이
GAP_X = 0.15         # 역끼리 간격
START_X = 0.4

for li, line in enumerate(lines):
    row_y = ROW_START_Y + li * ROW_H

    # 호선 라벨 (큰 원형 색칠)
    lbl = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(START_X), Inches(row_y),
        Inches(0.9), Inches(0.9)
    )
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = hex_to_rgb(line['color'])
    lbl.line.fill.background()
    ltf = lbl.text_frame
    ltf.margin_left = Emu(0); ltf.margin_right = Emu(0)
    ltf.margin_top = Emu(0); ltf.margin_bottom = Emu(0)
    ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = line['name']
    lr.font.name = '맑은 고딕'
    lr.font.size = Pt(13)
    lr.font.bold = True
    lr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 역 박스들
    for si, st in enumerate(line['stations']):
        x = START_X + 1.2 + si * (STATION_W + GAP_X)
        y = row_y + 0.15  # 호선 라벨 가운데 정렬
        sh = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(STATION_W), Inches(STATION_H)
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = hex_to_rgb(line['color'])
        sh.line.color.rgb = hex_to_rgb(line['color'])
        sh.line.width = Pt(1.5)

        stf = sh.text_frame
        stf.margin_left = Emu(36000); stf.margin_right = Emu(36000)
        stf.margin_top = Emu(18000); stf.margin_bottom = Emu(18000)
        stf.vertical_anchor = MSO_ANCHOR.MIDDLE
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        sr.text = st + '역'
        sr.font.name = '맑은 고딕'
        sr.font.size = Pt(13)
        sr.font.bold = True
        sr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# 안내 텍스트
note = slide3.shapes.add_textbox(Inches(0.4), Inches(ROW_START_Y + 5*ROW_H + 0.2),
                                  Inches(SLIDE_W_IN - 0.8), Inches(0.5))
ntf = note.text_frame
np_ = ntf.paragraphs[0]
nr = np_.add_run()
nr.text = "※ 환승역(잠실, 가락시장, 오금)은 여러 호선에 중복 표시되어 있어요"
nr.font.name = '맑은 고딕'
nr.font.size = Pt(11)
nr.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

out = '송파구/송파구_아파트_임장지도_v2.pptx'
prs.save(out)
print(f"완료! → {out}")
print(f"슬라이드 1: 지도 + 동별 라벨")
print(f"슬라이드 2: 가나다순 라벨 {len(apts_sorted)}개")
print(f"슬라이드 3: 송파구 지하철역 (호선별)")
