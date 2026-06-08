import sys, io, json, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from PIL import Image
import base64
from pptx import Presentation
from pptx.util import Inches

# ── 이미지 base64 ──────────────────────────────────────────
img = Image.open('송파구/송파구이미지2.png').convert('RGB')
IW, IH = img.size
small = img.resize((2000, int(IH * 2000 / IW)), Image.LANCZOS)
SW, SH = small.size
buf = io.BytesIO()
small.save(buf, format='JPEG', quality=88)
b64 = base64.b64encode(buf.getvalue()).decode()

# ── PPT 읽기 ───────────────────────────────────────────────
with open('송파구/프레젠테이션1_최종.pptx', 'rb') as f:
    prs = Presentation(f)
slide = prs.slides[0]
SLIDE_W = prs.slide_width   # EMU
SLIDE_H = prs.slide_height  # EMU

markers = []
for sh in list(slide.shapes):
    if not sh.has_text_frame: continue
    text = sh.text_frame.text.strip()
    if not text: continue

    # 슬라이드 내 위치 → 이미지 픽셀 위치 (박스 중심)
    cx_emu = sh.left + sh.width / 2
    cy_emu = sh.top + sh.height / 2
    px = cx_emu / SLIDE_W * SW
    py = cy_emu / SLIDE_H * SH

    # 배경색
    try:
        from pptx.enum.dml import MSO_THEME_COLOR
        fill = sh.fill
        if fill.type is not None:
            bg = '#{:06X}'.format(int(fill.fore_color.rgb))
        else:
            bg = 'transparent'
    except:
        bg = 'transparent'

    # 테두리색 → 글자색으로도 사용
    try:
        color = '#{:06X}'.format(int(sh.line.color.rgb))
    except:
        color = '#222222'

    # 글자색: 첫 번째 run에서 읽기
    font_color = '#111111'
    try:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.color and r.font.color.rgb:
                    font_color = '#{:06X}'.format(int(r.font.color.rgb))
                    break
    except:
        pass

    lines = [p.text.strip() for p in sh.text_frame.paragraphs if p.text.strip()]
    label = '<br>'.join(lines)

    markers.append({'x': round(px,1), 'y': round(py,1),
                    'label': label, 'color': color, 'bg': bg, 'fc': font_color})

print(f"마커: {len(markers)}개")
markers_js = json.dumps(markers, ensure_ascii=False)

# ── HTML ───────────────────────────────────────────────────
html = """<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1,maximum-scale=10,user-scalable=yes">
<title>송파구 임장지도</title>
<style>
*{margin:0;padding:0;box-sizing:border-box;}
body{background:#333;overflow:hidden;touch-action:none;}
#wrap{position:relative;display:inline-block;transform-origin:0 0;}
.lb{
  position:absolute;
  font-family:'Malgun Gothic','나눔고딕',sans-serif;
  font-size:11px;font-weight:bold;
  padding:3px 7px;border-radius:6px;border:2px solid;
  white-space:nowrap;line-height:1.55;
  box-shadow:1px 2px 5px rgba(0,0,0,0.3);
  transform:translate(-50%,-50%);
  cursor:default;
}
#legend{
  position:fixed;top:10px;left:10px;
  background:rgba(255,255,255,0.93);
  padding:8px 12px;border-radius:8px;
  border:1.5px solid #aaa;font-size:12px;
  font-family:'Malgun Gothic',sans-serif;
  line-height:1.8;z-index:999;
}
</style>
</head>
<body>
<div id="legend">
  <b>준공 연식</b><br>
  <span style="color:#C0392B">■</span> ~1989<br>
  <span style="color:#D35400">■</span> 1990~1999<br>
  <span style="color:#1E8449">■</span> 2000~2009<br>
  <span style="color:#1A5276">■</span> 2010~2019<br>
  <span style="color:#6C3483">■</span> 2020~
</div>
<div id="wrap">
  <img id="map-img" src="data:image/jpeg;base64,IMG_B64" draggable="false">
</div>
<script>
const RAW_W=SW_VAL, RAW_H=SH_VAL;
const markers=MARKERS_JS;
const wrap=document.getElementById('wrap');
let scale=1,offX=0,offY=0;
function init(){
  const s=Math.min(window.innerWidth/RAW_W,window.innerHeight/RAW_H);
  scale=s;offX=(window.innerWidth-RAW_W*s)/2;offY=(window.innerHeight-RAW_H*s)/2;apply();
}
function apply(){wrap.style.transform=`translate(${offX}px,${offY}px) scale(${scale})`;}

let touches={},lastDist=0,dragStart=null,startOff={x:0,y:0};
document.addEventListener('touchstart',e=>{
  e.preventDefault();
  [...e.changedTouches].forEach(t=>touches[t.identifier]={x:t.clientX,y:t.clientY});
  const pts=Object.values(touches);
  if(pts.length===1){dragStart={x:pts[0].x,y:pts[0].y};startOff={x:offX,y:offY};}
  else if(pts.length===2){lastDist=Math.hypot(pts[1].x-pts[0].x,pts[1].y-pts[0].y);}
},{passive:false});
document.addEventListener('touchmove',e=>{
  e.preventDefault();
  [...e.changedTouches].forEach(t=>touches[t.identifier]={x:t.clientX,y:t.clientY});
  const pts=Object.values(touches);
  if(pts.length===1&&dragStart){offX=startOff.x+(pts[0].x-dragStart.x);offY=startOff.y+(pts[0].y-dragStart.y);}
  else if(pts.length===2){
    const dist=Math.hypot(pts[1].x-pts[0].x,pts[1].y-pts[0].y);
    const mid={x:(pts[0].x+pts[1].x)/2,y:(pts[0].y+pts[1].y)/2};
    const nx=(mid.x-offX)/scale,ny=(mid.y-offY)/scale;
    scale*=dist/lastDist;scale=Math.max(0.3,Math.min(scale,15));
    offX=mid.x-nx*scale;offY=mid.y-ny*scale;lastDist=dist;
  }
  apply();
},{passive:false});
document.addEventListener('touchend',e=>{
  [...e.changedTouches].forEach(t=>delete touches[t.identifier]);
  if(!Object.keys(touches).length)dragStart=null;
});
let drag=false,mx,my,ox,oy;
wrap.addEventListener('mousedown',e=>{drag=true;mx=e.clientX;my=e.clientY;ox=offX;oy=offY;});
window.addEventListener('mousemove',e=>{if(!drag)return;offX=ox+(e.clientX-mx);offY=oy+(e.clientY-my);apply();});
window.addEventListener('mouseup',()=>drag=false);
window.addEventListener('wheel',e=>{
  e.preventDefault();const f=e.deltaY<0?1.12:0.89;
  const nx=(e.clientX-offX)/scale,ny=(e.clientY-offY)/scale;
  scale*=f;scale=Math.max(0.3,Math.min(scale,15));
  offX=e.clientX-nx*scale;offY=e.clientY-ny*scale;apply();
},{passive:false});

markers.forEach(m=>{
  const el=document.createElement('div');
  el.className='lb';el.innerHTML=m.label;
  el.style.left=m.x+'px';el.style.top=m.y+'px';
  el.style.color=m.fc||m.color;el.style.background=m.bg;el.style.borderColor=m.color;
  if(m.bg==='transparent'){el.style.border='none';el.style.boxShadow='none';el.style.padding='1px 3px';}
  wrap.appendChild(el);
});
init();
window.addEventListener('resize',init);
</script>
</body>
</html>"""

html = (html
    .replace('IMG_B64', b64)
    .replace('SW_VAL', str(SW))
    .replace('SH_VAL', str(SH))
    .replace('MARKERS_JS', markers_js))

out = '송파구/송파구_임장지도_모바일.html'
with open(out, 'w', encoding='utf-8') as f:
    f.write(html)
print(f"저장: {out}  ({len(html.encode())//1024}KB)")
