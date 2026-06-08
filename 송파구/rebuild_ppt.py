import sys, io, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
import pandas as pd
NUM_RE = re.compile(r'\d+')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 1. 평형 데이터 빌드 ──────────────────────────────────────
df = pd.read_excel('송파구/송파구 루시퍼홍.xlsx', engine='openpyxl')
songpa = df[df['시군구']=='송파구']

apt_pyeong = {}
for code, group in songpa.groupby('단지코드'):
    name  = group['단지명'].iloc[0]
    yr2   = str(group['사용승인일'].iloc[0])[2:4]
    total = int(group['총세대수'].iloc[0])
    p_dict, big = {}, 0
    for _, row in group.iterrows():
        sqm   = row['전용면적']
        units = int(row['평형별 세대수'])
        if sqm >= 85:
            big += units
        else:
            pyeong_name = str(row['평형명'])
            nums = NUM_RE.findall(pyeong_name)
            if nums:
                supply_sqm = int(nums[0])
                p = int(supply_sqm / 3.3058 + 0.5)
            else:
                p = int(sqm / 3.3058 + 0.5)
            p_dict[p] = p_dict.get(p, 0) + units
    lines = [f"{p}p({p_dict[p]:,}^)" for p in sorted(p_dict)]
    if big: lines.append(f"큰평형({big:,}^)")
    apt_pyeong[(name, yr2, total)] = lines

# ── 2. PPT 열기 ───────────────────────────────────────────────
prs = Presentation('송파구/프레젠테이션1_색상적용_초등학교포함.pptx')
slide  = prs.slides[0]
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

pattern = re.compile(r"(\d{2})\s*[''']\s*([\d,]+)\s*\^")

# ── 3. 기존 아파트 박스: 평형 추가 + 높이 자동 조정 ──────────
LINE_H_EMU = Pt(8.5) * 12700 // 100   # 약 1줄 높이
PER_LINE   = Emu(75000)                # 줄당 추가 높이 (2줄박스 150332/2 기준)

def ref_run(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            if r.text.strip(): return r
    return None

updated = 0
for sh in list(slide.shapes):
    if not sh.has_text_frame: continue
    text = sh.text_frame.text.strip()
    m = pattern.search(text)
    if not m: continue

    yr2   = m.group(1)
    total = int(m.group(2).replace(',',''))
    name_part = text[:m.start()].replace('\n',' ').replace('|',' ').strip()
    while '  ' in name_part: name_part = name_part.replace('  ',' ')

    # 매칭
    best = None
    for (n, y, t), lines in apt_pyeong.items():
        if y==yr2 and t==total and (n==name_part or n in name_part or name_part in n):
            best = (n, lines); break
    if not best:
        for (n, y, t), lines in apt_pyeong.items():
            if y==yr2 and t==total:
                best = (n, lines); break
    if not best: continue

    tf   = sh.text_frame
    rrun = ref_run(tf)

    for line in best[1]:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        if rrun:
            if rrun.font.name:  r.font.name  = rrun.font.name
            if rrun.font.size:  r.font.size  = rrun.font.size
            if rrun.font.bold is not None: r.font.bold = rrun.font.bold
            try:
                if rrun.font.color.rgb: r.font.color.rgb = rrun.font.color.rgb
            except: pass

    # 높이 조정: 원래 높이 + 추가된 줄 수 * 줄당 높이
    n_added = len(best[1])
    sh.height = sh.height + PER_LINE * n_added
    updated += 1

print(f"기존 박스 업데이트: {updated}개")

# ── 4. 호수임광 박스 추가 ──────────────────────────────────────
hsu_lines = apt_pyeong.get(('호수임광','95',227), [])
hsu_all   = ['호수임광', "95' 227^"] + hsu_lines

# 연식 색상 (95년 → 1990년대 → 주황)
fill_hex, line_hex = 'FEF0E0', 'D35400'
def hex_rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

BOX_W = Inches(1.15)
n_new = len(hsu_all)
BOX_H = PER_LINE * n_new

# 오른쪽 위
left = SLIDE_W - BOX_W - Emu(50000)
top  = Emu(60000)

sh_new = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, BOX_W, BOX_H)
sh_new.fill.solid()
sh_new.fill.fore_color.rgb = hex_rgb(fill_hex)
sh_new.line.color.rgb      = hex_rgb(line_hex)
sh_new.line.width          = Pt(1.5)

tf = sh_new.text_frame
tf.margin_left   = Emu(36000)
tf.margin_right  = Emu(36000)
tf.margin_top    = Emu(16000)
tf.margin_bottom = Emu(16000)
tf.word_wrap     = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

for i, line in enumerate(hsu_all):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = line
    r.font.name = '맑은 고딕'
    r.font.size = Pt(8)
    r.font.bold = (i==0)
    r.font.color.rgb = RGBColor(0x11,0x11,0x11) if i<2 else hex_rgb(line_hex)

print(f"호수임광 추가: {hsu_all}")

# ── 5. 저장 ──────────────────────────────────────────────────
out = '송파구/프레젠테이션1_최종.pptx'
prs.save(out)
print(f"저장: {out}")
