import sys, io, re
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
import pandas as pd
from pptx import Presentation
from copy import deepcopy

# 1) Excel 평형 데이터 빌드
df = pd.read_excel('송파구/송파구 루시퍼홍.xlsx', engine='openpyxl')
songpa = df[df['시군구']=='송파구']

apt_pyeong = {}  # (단지명, 연도2자리, 총세대수) → [평형 라인들]

for code, group in songpa.groupby('단지코드'):
    name = group['단지명'].iloc[0]
    yr = str(group['사용승인일'].iloc[0])[:4]
    yr2 = yr[2:]
    total = int(group['총세대수'].iloc[0])

    p_dict = {}
    big = 0
    for _, row in group.iterrows():
        sqm = row['전용면적']
        units = int(row['평형별 세대수'])
        if sqm >= 86:
            big += units
        elif sqm == 84 or sqm == 85:
            p_dict[25] = p_dict.get(25, 0) + units
        else:
            p = int(sqm / 3.3058 + 0.5)
            p_dict[p] = p_dict.get(p, 0) + units

    lines = []
    for p in sorted(p_dict.keys()):
        lines.append(f"{p}p({p_dict[p]:,}^)")
    if big > 0:
        lines.append(f"큰평형({big:,}^)")

    apt_pyeong[(name, yr2, total)] = lines

print(f"단지: {len(apt_pyeong)}개")

# 2) PPT 업데이트
prs = Presentation('송파구/프레젠테이션1_색상적용_초등학교포함.pptx')
slide = prs.slides[0]

pattern = re.compile(r"(\d{2})\s*[''']\s*([\d,]+)\s*\^")

updated = 0
not_found = []
for sh in list(slide.shapes):
    if not sh.has_text_frame: continue
    text = sh.text_frame.text.strip()
    m = pattern.search(text)
    if not m: continue

    yr2 = m.group(1)
    total = int(m.group(2).replace(',', ''))

    # 이름 추출
    name_part = text[:m.start()].strip()
    name_part = name_part.replace('\n', ' ').replace('|', ' ').strip()
    while '  ' in name_part:
        name_part = name_part.replace('  ', ' ')

    # 매칭
    best = None
    # 1차: 이름+연도+세대수
    for (n, y, t), lines in apt_pyeong.items():
        if y == yr2 and t == total:
            if n == name_part or n in name_part or name_part in n:
                best = (n, lines)
                break
    # 2차: 연도+세대수만
    if not best:
        for (n, y, t), lines in apt_pyeong.items():
            if y == yr2 and t == total:
                best = (n, lines)
                break

    if not best:
        not_found.append(text[:50])
        continue

    # 기존 마지막 run의 스타일 복사용
    tf = sh.text_frame
    ref_run = None
    for p in tf.paragraphs:
        for r in p.runs:
            if r.text.strip():
                ref_run = r

    # 평형 라인 추가
    for line in best[1]:
        p = tf.add_paragraph()
        new_r = p.add_run()
        new_r.text = line
        if ref_run is not None:
            if ref_run.font.name:
                new_r.font.name = ref_run.font.name
            if ref_run.font.size:
                new_r.font.size = ref_run.font.size
            if ref_run.font.bold is not None:
                new_r.font.bold = ref_run.font.bold
            try:
                if ref_run.font.color and ref_run.font.color.rgb:
                    new_r.font.color.rgb = ref_run.font.color.rgb
            except:
                pass

    updated += 1

print(f"업데이트: {updated}개")
if not_found:
    print(f"매칭 실패 ({len(not_found)}개):")
    for t in not_found:
        print(f"  - {t}")

prs.save('송파구/프레젠테이션1_평형추가.pptx')
print("저장: 송파구/프레젠테이션1_평형추가.pptx")
