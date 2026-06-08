
"""
슬라이드14 오른쪽 표 채우기
- 연식 기준으로 단지명+연식 / 84p 기준 가격 입력
"""
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from lxml import etree

A = 'http://schemas.openxmlformats.org/drawingml/2006/main'

PPT_PATH = 'C:/Temp/songpa_work.pptx'
prs = Presentation(PPT_PATH)
slide = prs.slides[13]
shapes = list(slide.shapes)

# ── 파싱 함수 ─────────────────────────────────────────────────
HDR_PAT   = re.compile(r'^(.+?)\s+(\d{2})\s+[\d,]+\^', re.DOTALL)
PRICE_PAT = re.compile(r'^(\d+)p\s+([\d.]+)/([\d.]+|-)')

def parse_header(text):
    text2 = text.replace('\n', ' ').strip()
    m = HDR_PAT.match(text2)
    if not m: return None, None
    name = m.group(1).strip()
    yr2  = int(m.group(2))
    year = 2000 + yr2 if yr2 <= 30 else 1900 + yr2
    return name, year

def get_best_price(content_text, target=84):
    best_diff, best_str = float('inf'), None
    for line in content_text.split('\n'):
        m = PRICE_PAT.match(line.strip())
        if m:
            p = int(m.group(1))
            diff = abs(p - target)
            if diff < best_diff:
                best_diff = diff
                best_str = f'{p}p {m.group(2)}/{m.group(3)}'
    return best_str

# ── 헤더/가격 박스 탐지 ───────────────────────────────────────
header_shapes = {}   # id(sh) → (sh, name, year)
content_shapes = {}  # id(sh) → (sh, text)

for sh in shapes:
    if not sh.has_text_frame: continue
    text = sh.text_frame.text.strip()
    if not text: continue
    name, year = parse_header(text)
    if name and year:
        header_shapes[id(sh)] = (sh, name, year)
    elif PRICE_PAT.search(text) or '가격미확인' in text or '가격 미확인' in text:
        content_shapes[id(sh)] = (sh, text)

print(f'헤더 박스 {len(header_shapes)}개, 가격 박스 {len(content_shapes)}개')

# ── 헤더↔컨텐츠 매칭 (가장 가까운 위쪽 헤더) ─────────────────
apartments = []

for cid, (csh, ctext) in content_shapes.items():
    best_hsh, best_dist = None, float('inf')
    for hid, (hsh, hname, hyear) in header_shapes.items():
        dy = csh.top - hsh.top          # 양수 = 컨텐츠가 헤더 아래
        dx = abs(csh.left - hsh.left)
        if dy > 0 and dx < int(Pt(120)):
            dist = dy + dx * 0.3
            if dist < best_dist:
                best_dist = dist
                best_hsh = (hsh, hname, hyear)

    if best_hsh:
        hsh, name, year = best_hsh
        price_str = get_best_price(ctext) or '가격미확인'
        apartments.append((name, year, price_str))
        print(f'  {name} {year}: {price_str}')

# ── 연식별 분류 ───────────────────────────────────────────────
ERA_RANGES = [
    (2019, 2099),   # row1 구신축
    (2014, 2018),   # row2 해신축
    (2004, 2013),   # row3 준구축
    (0,    2003),   # row4 구축
]
era_apts = [[] for _ in range(4)]
for name, year, price_str in apartments:
    for i, (y_from, y_to) in enumerate(ERA_RANGES):
        if y_from <= year <= y_to:
            era_apts[i].append((name, year, price_str))
            break

def sort_key(apt):
    name, year, price = apt
    # 가격 문자열에서 매매가 숫자 추출 (예: "84p 33.5/15.5" → 33.5)
    m = re.search(r'p\s+([\d.]+)/', price)
    mae = float(m.group(1)) if m else 0.0
    return (year, -mae)   # 연식 오름차순, 가격 내림차순

for i in range(len(era_apts)):
    era_apts[i].sort(key=sort_key)

print('\n=== 연식별 단지 (정렬 후) ===')
for i, apts in enumerate(era_apts):
    for name, year, price in apts:
        yr2 = str(year)[-2:]
        print(f'  row{i+1} | {name}{yr2}\' | {price}')

# ── 표 셀 채우기 ──────────────────────────────────────────────
table_sh = shapes[6]
tbl = table_sh.table
FONT_SIZE = Pt(10)

def get_txBody(cell):
    tc = cell._tc
    return tc.find(f'{{{A}}}txBody')

def clear_cell(cell):
    txBody = get_txBody(cell)
    if txBody is None: return
    for p in list(txBody.findall(f'{{{A}}}p')):
        txBody.remove(p)

def add_para(cell, text, font_size=FONT_SIZE, bold=False, color='111111'):
    txBody = get_txBody(cell)
    if txBody is None: return
    p_el = etree.SubElement(txBody, f'{{{A}}}p')
    pPr = etree.SubElement(p_el, f'{{{A}}}pPr')
    pPr.set('algn', 'l')
    # 줄간격 100%, 앞뒤 간격 0
    lnSpc = etree.SubElement(pPr, f'{{{A}}}lnSpc')
    etree.SubElement(lnSpc, f'{{{A}}}spcPct').set('val', '100000')
    spcBef = etree.SubElement(pPr, f'{{{A}}}spcBef')
    etree.SubElement(spcBef, f'{{{A}}}spcPts').set('val', '0')
    spcAft = etree.SubElement(pPr, f'{{{A}}}spcAft')
    etree.SubElement(spcAft, f'{{{A}}}spcPts').set('val', '0')
    # 런
    r_el = etree.SubElement(p_el, f'{{{A}}}r')
    rPr = etree.SubElement(r_el, f'{{{A}}}rPr')
    rPr.set('lang', 'ko-KR')
    rPr.set('sz', str(int(font_size.pt * 100)))
    rPr.set('b', '1' if bold else '0')
    rPr.set('dirty', '0')
    sf = etree.SubElement(rPr, f'{{{A}}}solidFill')
    etree.SubElement(sf, f'{{{A}}}srgbClr').set('val', color)
    etree.SubElement(r_el, f'{{{A}}}t').text = text

for row_idx, apts in enumerate(era_apts):
    row_num = row_idx + 1
    dcell = tbl.cell(row_num, 1)
    pcell = tbl.cell(row_num, 2)
    clear_cell(dcell)
    clear_cell(pcell)

    if not apts:
        add_para(dcell, '-', FONT_SIZE, color='888888')
        add_para(pcell, '-', FONT_SIZE, color='888888')
        continue

    for name, year, price in apts:
        yr2   = str(year)[-2:]
        label = f"{name}{yr2}'"
        add_para(dcell, label, FONT_SIZE)
        add_para(pcell, price, FONT_SIZE)

    print(f'표 row{row_num} 채움: {len(apts)}개')

prs.save(PPT_PATH)
print('\n저장 완료:', PPT_PATH)
