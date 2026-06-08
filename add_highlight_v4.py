
"""
형광펜 스타일 가격박스 v4
1. 테두리 1pt
2. 기존 코드 텍스트박스 삭제
3. 위치 정렬 (슬라이드 영역 초과 방지)
4. 텍스트 길이에 따라 가로폭 자동 조절
5. 큰평형/재건축 등 추가정보: 흰색 형광펜
6. 형광펜 색 대비에 따라 글자색 자동 결정
7. 갤러리아팰리스, 레이크팰리스 동일 형식 추가
"""
import re, json, openpyxl
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
SLIDE_W = 9144000   # EMU (10인치 슬라이드)
SLIDE_H = 6858000   # EMU (7.5인치 슬라이드)

RAINBOW = {
    35:('000000','FFFFFF'), 34:('1A1A1A','FFFFFF'), 33:('2B2B2B','FFFFFF'),
    32:('424242','FFFFFF'), 31:('616161','FFFFFF'), 30:('757575','FFFFFF'),
    29:('4A0000','FFFFFF'), 28:('6A0F0F','FFFFFF'), 27:('8D1515','FFFFFF'),
    26:('C62828','FFFFFF'), 25:('E53935','FFFFFF'), 24:('F4511E','FFFFFF'),
    23:('FB8C00','FFFFFF'), 22:('FFB300','FFFFFF'), 21:('F9E400','000000'),
    20:('9CCC65','000000'), 19:('43A047','FFFFFF'), 18:('00897B','FFFFFF'),
    17:('00ACC1','000000'), 16:('29B6F6','000000'), 15:('1E88E5','FFFFFF'),
    14:('1565C0','FFFFFF'), 13:('1A237E','FFFFFF'), 12:('7B2FBE','FFFFFF'),
    11:('B07FE0','FFFFFF'), 10:('D7B8F5','000000'),
}

def fix(v):
    if v is None or isinstance(v, str): return None
    return v // 10 if v > 1000000 else v

def num(m):
    if m is None: return '-'
    v = m / 10000
    s = f'{v:.1f}'
    return s[:-2] if s.endswith('.0') else s

def get_color(mae):
    if mae is None: return ('595959','FFFFFF')
    eok = int(mae // 10000)
    if eok >= 35: return RAINBOW[35]
    return RAINBOW.get(eok, ('595959','FFFFFF'))

def yr2(y):
    if not y: return '??'
    return str(int(y))[-2:]

def sd_str(s):
    if s is None: return '?'
    return str(int(s)) if isinstance(s, float) else str(s)

def text_width_pt(text, font_pt=8):
    """텍스트 길이 추정 (pt 단위)"""
    w = 0
    for c in text:
        if '가' <= c <= '힣':   # 한글
            w += font_pt * 1.0
        elif '一' <= c <= '鿿': # 한자
            w += font_pt * 1.0
        elif c == ' ':
            w += font_pt * 0.3
        elif c in '()':
            w += font_pt * 0.4
        elif c in '/.,':
            w += font_pt * 0.35
        else:
            w += font_pt * 0.6
    return w

def calc_box_w(all_lines, font_pt=8, min_pt=55, padding_pt=10):
    """모든 라인 중 가장 긴 것 기준으로 폭 계산"""
    max_w = 0
    for line in all_lines:
        max_w = max(max_w, text_width_pt(line, font_pt))
    total = max_w + padding_pt
    return max(int(Pt(total)), int(Pt(min_pt)))

def set_highlight(run, bg_hex, fg_hex):
    rPr = run._r.find(f'{{{A}}}rPr')
    if rPr is None:
        rPr = etree.SubElement(run._r, f'{{{A}}}rPr')
    for tag in [f'{{{A}}}highlight', f'{{{A}}}solidFill']:
        for el in rPr.findall(tag):
            rPr.remove(el)
    # DrawingML 스키마 순서: solidFill(글자색) → highlight(배경색) 순서여야 적용됨
    sf = etree.SubElement(rPr, f'{{{A}}}solidFill')
    clr2 = etree.SubElement(sf, f'{{{A}}}srgbClr')
    clr2.set('val', fg_hex)
    hl = etree.SubElement(rPr, f'{{{A}}}highlight')
    clr = etree.SubElement(hl, f'{{{A}}}srgbClr')
    clr.set('val', bg_hex)

def set_no_fill(shape):
    sp = shape._element
    spPr = sp.find(f'{{{P_NS}}}spPr')
    if spPr is None: spPr = sp.find(f'{{{A}}}spPr')
    if spPr is None: return
    for tag in [f'{{{A}}}solidFill', f'{{{A}}}noFill', f'{{{A}}}gradFill']:
        for el in spPr.findall(tag):
            spPr.remove(el)
    etree.SubElement(spPr, f'{{{A}}}noFill')

def set_border(shape, hex_color='333333', width_pt=1.0):
    sp = shape._element
    spPr = sp.find(f'{{{P_NS}}}spPr')
    if spPr is None: spPr = sp.find(f'{{{A}}}spPr')
    if spPr is None: return
    for old_ln in spPr.findall(f'{{{A}}}ln'):
        spPr.remove(old_ln)
    ln = etree.SubElement(spPr, f'{{{A}}}ln')
    ln.set('w', str(int(Pt(width_pt))))   # 1pt = 12700 EMU
    sf = etree.SubElement(ln, f'{{{A}}}solidFill')
    clr = etree.SubElement(sf, f'{{{A}}}srgbClr')
    clr.set('val', hex_color)

def set_bPr(tf, l=36000, r=36000, t=18000, b=18000, anchor='t'):
    bPr = tf._txBody.find(f'{{{A}}}bodyPr')
    if bPr is not None:
        bPr.set('lIns', str(l)); bPr.set('rIns', str(r))
        bPr.set('tIns', str(t)); bPr.set('bIns', str(b))
        bPr.set('anchor', anchor)

# ── hanpan5.html에서 추가 정보 파싱 ─────────────────────────────
def load_extra_info(html_path):
    with open(html_path, 'rb') as f:
        content = f.read().decode('utf-8')
    m = re.search(r'apartments=(\[.*?\]);', content, re.DOTALL)
    if not m: return {}
    arr = json.loads(m.group(1))
    price_pat = re.compile(r'^\d+p\(')
    result = {}
    for apt in arr:
        name_raw = apt.get('name', '').split('\n')[0].strip()
        text = apt.get('text', '')
        lines = text.split('\n')
        extra = [l.strip() for l in lines[1:] if l.strip() and not price_pat.match(l.strip())]
        if extra:
            result[name_raw] = extra
    return result

extra_db = load_extra_info('C:/Temp/hanpan5.html')

SEARCH_KEYS = {
    '잠실엘스': '잠실엘스',
    '리센츠': '리센츠',
    '잠실주공5단지': '잠실주공5단지',
    '트리지움': '트리지움',
    '우성1,2,3차': '우성1,2,3차',
    '아시아선수촌': '아시아선수촌',
    '우성4차': '우성4차',
    '갤러리아팰리스': '갤러리아팰리스',
    '레이크팰리스': '레이크팰리스',
}

def get_extra_lines(apt_name):
    for key in SEARCH_KEYS:
        if key in apt_name:
            for html_name, lines in extra_db.items():
                if key in html_name:
                    return lines
    return []

# ── Excel 로드 ─────────────────────────────────────────────────
wb = openpyxl.load_workbook('C:/Temp/sisae.xlsx', data_only=True)
ws = wb['송파구']
CODE_TARGETS = [22627, 22746, 639, 19127, 12861, 635, 637, 640]
EXTRA_CODES  = [12236, 15011]   # 갤러리아팰리스, 레이크팰리스
ALL_TARGETS  = CODE_TARGETS + EXTRA_CODES

data = {}
for row in range(6, ws.max_row + 1):
    code = ws.cell(row, 4).value
    if code in ALL_TARGETS:
        data.setdefault(code, []).append({
            'name':   ws.cell(row, 6).value,
            'pyeong': ws.cell(row, 7).value,
            'mae':    ws.cell(row, 8).value,
            'jun':    ws.cell(row, 9).value,
            'sedae':  ws.cell(row, 13).value,
            'year':   ws.cell(row, 12).value,
        })

# ── PPT 로드 ───────────────────────────────────────────────────
prs = Presentation('C:/Temp/songpa_latest2.pptx')
slide = prs.slides[13]

# ── 1단계: 위치 기록 & 삭제 대상 수집 ─────────────────────────
code_positions = {}    # code_int → (left, top, width, height)
galrae_positions = {}  # code → (center_x, top)
delete_elems = []      # 삭제할 XML 요소 목록 (id() 대신 요소 직접 참조)

for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    text = sh.text_frame.text.strip()

    # 코드 텍스트박스
    try:
        val = int(text)
        if val in CODE_TARGETS:
            code_positions[val] = (sh.left, sh.top, sh.width, sh.height)
            delete_elems.append(sh._element)
            continue
    except:
        pass

    # 갤러리아팰리스 기존 박스
    if '갤러리아팰리스' in text:
        if 12236 not in galrae_positions:
            galrae_positions[12236] = (sh.left + sh.width // 2, sh.top)
        delete_elems.append(sh._element)

    # 레이크팰리스 기존 박스
    elif '레이크팰리스' in text:
        if 15011 not in galrae_positions:
            galrae_positions[15011] = (sh.left + sh.width // 2, sh.top)
        delete_elems.append(sh._element)

# 갤/레 헤더 근처 다른 박스도 삭제 (XML 요소 기준으로 중복 체크)
delete_elem_set = set(id(el) for el in delete_elems)
for sh in slide.shapes:
    if id(sh._element) in delete_elem_set:
        continue
    if not sh.has_text_frame:
        continue
    for code, (cx, ht) in galrae_positions.items():
        hdr_l = cx - int(Pt(110))
        hdr_r = cx + int(Pt(110))
        if hdr_l <= sh.left <= hdr_r and ht <= sh.top <= ht + int(Pt(400)):
            delete_elems.append(sh._element)
            delete_elem_set.add(id(sh._element))
            break

# ── 2단계: 삭제 실행 ──────────────────────────────────────────
for el in delete_elems:
    parent = el.getparent()
    if parent is not None:
        parent.remove(el)
print(f'삭제 완료: {len(delete_elems)}개')

# ── 3단계: 박스 추가 ──────────────────────────────────────────
HDR_H      = int(Pt(12))    # 헤더 높이
ROW_H      = int(Pt(11))    # 가격행 높이
ROW_H_EX   = int(Pt(10))    # 추가정보행 높이
T_INS      = 10000          # 내부 상하 여백 (EMU, ≈0.8pt)
FONT       = Pt(8)

def add_boxes(code, cx, box_top, rows):
    """단지 박스 추가 (cx: 박스 중앙 x, box_top: 헤더 상단 y)"""
    r0 = rows[0]
    name  = r0['name'] or f'코드{code}'
    year  = yr2(r0['year'])
    sedae = sd_str(r0['sedae'])
    hdr_text = f'{name} {year} {sedae}^'

    # 가격행
    price_lines = []
    for r in rows:
        mae_raw = r['mae'] if not isinstance(r['mae'], str) else None
        jun_raw = r['jun'] if not isinstance(r['jun'], str) else None
        mae = fix(mae_raw)
        jun = fix(jun_raw)
        if mae is None: continue
        txt = f'{r["pyeong"]}p {num(mae)}/{num(jun)} ({num(mae - jun)})' if jun else f'{r["pyeong"]}p {num(mae)}/-'
        bg, fg = get_color(mae)
        price_lines.append((txt, bg, fg))

    if not price_lines:
        price_lines = [('가격 미확인', '595959', 'FFFFFF')]

    # 추가 정보행 (큰평형, 용적률, 재건축)
    extra_lines = get_extra_lines(name)

    # 가로폭 계산
    all_texts = [hdr_text] + [t for t, _, _ in price_lines] + extra_lines
    BOX_W = calc_box_w(all_texts, font_pt=8, min_pt=55, padding_pt=12)

    # 위치 보정 (슬라이드 벗어나지 않도록)
    L = cx - BOX_W // 2
    if L < int(Pt(3)): L = int(Pt(3))
    if L + BOX_W > SLIDE_W - int(Pt(3)):
        L = SLIDE_W - BOX_W - int(Pt(3))

    n_price = len(price_lines)
    n_extra = len(extra_lines)
    # 세로 높이: 행 수 × 행 높이 + 상하 내부 여백
    content_h = n_price * ROW_H + n_extra * ROW_H_EX + T_INS * 2
    total_h   = HDR_H + content_h

    T = box_top
    if T < int(Pt(3)): T = int(Pt(3))
    if T + total_h > SLIDE_H - int(Pt(3)):
        T = SLIDE_H - total_h - int(Pt(3))

    # ── 헤더 박스 ─────────────────────────────────────────────
    hdr_box = slide.shapes.add_textbox(L, T, BOX_W, HDR_H)
    hdr_box.fill.solid()
    hdr_box.fill.fore_color.rgb = RGBColor.from_string('111111')
    tf = hdr_box.text_frame
    tf.word_wrap = False
    set_bPr(tf, l=36000, r=36000, t=0, b=0, anchor='ctr')
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = hdr_text
    run.font.size = FONT
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string('FFFFFF')

    # ── 가격+추가정보 박스 ────────────────────────────────────
    price_box = slide.shapes.add_textbox(L, T + HDR_H, BOX_W, content_h)
    set_no_fill(price_box)
    set_border(price_box, '333333', 1.0)

    tf2 = price_box.text_frame
    tf2.word_wrap = False
    set_bPr(tf2, l=36000, r=36000, t=T_INS, b=T_INS, anchor='t')

    # 가격행 (형광펜)
    for i, (txt, bg, fg) in enumerate(price_lines):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run2 = para.add_run()
        run2.text = txt
        run2.font.size = FONT
        run2.font.bold = True
        set_highlight(run2, bg, fg)

    # 추가 정보행 (흰색 형광펜, 진회색 글씨)
    for extra in extra_lines:
        para = tf2.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run3 = para.add_run()
        run3.text = extra
        run3.font.size = Pt(7.5)
        run3.font.bold = False
        set_highlight(run3, 'FFFFFF', '333333')

    print(f'[{code}] {name}: {n_price}가격+{n_extra}추가 / 폭={BOX_W//12700:.0f}pt 높이={content_h//12700:.0f}pt')

# ── 코드 기반 8개 단지 ─────────────────────────────────────────
for code in CODE_TARGETS:
    if code not in code_positions:
        print(f'[{code}] 슬라이드에 없음')
        continue
    rows = data.get(code, [])
    if not rows:
        print(f'[{code}] Excel 없음')
        continue

    sl, st, sw, sh_h = code_positions[code]
    cx = sl + sw // 2

    rows_valid = [r for r in rows if fix(r['mae'] if not isinstance(r['mae'], str) else None) is not None]
    n_rows_est = max(len(rows_valid), 1)
    extra_est  = len(get_extra_lines(rows[0]['name'] or ''))
    total_h_est = HDR_H + n_rows_est * ROW_H + extra_est * ROW_H_EX + T_INS * 2
    T = st - total_h_est - int(Pt(4))
    if T < int(Pt(3)):
        T = st + sh_h + int(Pt(4))

    add_boxes(code, cx, T, rows)

# ── 갤러리아팰리스 & 레이크팰리스 ────────────────────────────
for code in EXTRA_CODES:
    if code not in galrae_positions:
        print(f'[{code}] 갤/레 위치 못 찾음')
        continue
    rows = data.get(code, [])
    if not rows:
        print(f'[{code}] Excel 없음')
        continue

    cx, top = galrae_positions[code]
    add_boxes(code, cx, top, rows)

prs.save('C:/Temp/songpa_highlight_v4.pptx')
print('\n저장: C:/Temp/songpa_highlight_v4.pptx')
